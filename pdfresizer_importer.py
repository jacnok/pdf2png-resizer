import os
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from io import BytesIO
import moviepy.editor as mp

def import_pdf(pdf_path):
    pdf_document = fitz.open(pdf_path)
    images = []
    
    for page_num in range(pdf_document.page_count):
        page = pdf_document.load_page(page_num)
        pix = page.get_pixmap()
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        images.append(img)
        
    return images

def import_images(image_folder, extensions=[".png", ".jpg", ".jpeg"]):
    images = []
    
    for filename in os.listdir(image_folder):
        if any(filename.lower().endswith(ext) for ext in extensions):
            img_path = os.path.join(image_folder, filename)
            img = Image.open(img_path)
            images.append(img)
            
    return images

def import_pptx(pptx_path, output_folder):
    # TODO: get this to actually work, because we're having issues with the pptx and group shapes
    
    prs = Presentation(pptx_path)
    images = []
    media_files = []
    
    for slide_num, slide in enumerate(prs.slides):
        slide_image = slide_to_image(slide, slide_num, output_folder)
        if slide_image:
            images.append(slide_image)
        slide_media = extract_media(slide, slide_num, output_folder)
        media_files.extend(slide_media)
    
    return images, media_files

def slide_to_image(slide, slide_num, output_folder):
    # Create an image of the slide
    img_path = os.path.join(output_folder, f"slide_{slide_num + 1}.png")
    slide.shapes._spTree.write_image(img_path)
    return Image.open(img_path)

def extract_media(slide, slide_num, output_folder):
    media_files = []
    
    for shape in slide.shapes:
        if not shape.has_media:
            continue
        
        media_part = shape.part.related_parts[shape._element.blip_rId]
        media_extension = media_part.content_type.split('/')[-1]
        media_filename = f"slide_{slide_num + 1}_media_{len(media_files) + 1}.{media_extension}"
        media_path = os.path.join(output_folder, media_filename)
        
        with open(media_path, 'wb') as f:
            f.write(media_part.blob)
        
        media_files.append(media_path)
    
    return media_files

def create_mp4_from_media(media_path, output_path):
    clip = mp.VideoFileClip(media_path)
    clip.write_videofile(output_path, codec="libx264", audio_codec="aac")
